import os
from datetime import datetime
from pptx import Presentation
import unicodedata
import glob

"""
extract_text_from_slide_decks_by_slide
extract text from a .pptx slide deck
ideally,
extract the text from each slide,

"""


def remove_control_chars(s):
    # Create a translation table that maps all control characters to None
    control_chars = dict.fromkeys(range(0, 32), None)
    # Add the U+000B character explicitly
    control_chars[11] = None
    # Translate the string using the translation table
    return s.translate(control_chars)


def extract_text_to_file(pptx_file):
    # Get the PowerPoint file name without extension
    file_name = os.path.splitext(os.path.basename(pptx_file))[0]

    # Specify the directory path you want to create
    save_here_directory_path = f"pptx_chunks/{file_name}_chunks"

    # Check if the directory exists
    if not os.path.exists(save_here_directory_path):
        # If the directory doesn't exist, create it
        os.makedirs(save_here_directory_path)
        print(f"Directory '{save_here_directory_path}' created successfully.")
    else:
        print(f"Directory '{save_here_directory_path}' already exists.")

    # Open the PowerPoint file
    prs = Presentation(pptx_file)

    # Get the current timestamp
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    slide_number = 1

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # f.write(shape.text + "\n")

                raw_text = shape.text
                utf8_text = remove_control_chars(raw_text)

                # save slide text
                slide_text.append(utf8_text)

        # if not empty
        if slide_text:
            
            # # Inspection
            # print("slide_text")
            # print(slide_text)

            # Create the output file name with slide number and timestamp
            output_file = (
                f"{save_here_directory_path}/{file_name}_{slide_number}_{timestamp}.txt"
            )

            # Write slide text to the output file
            with open(output_file, "w", encoding="utf-8") as f:
                f.write("\n".join(slide_text))

            print(f"Slide {slide_number} text saved to {output_file}")

        slide_number += 1

    print(f"Text extracted and saved to {output_file}")


if __name__ == "__main__":
    # Get the list of .pptx files in the current working directory

    directory_to_process = "files_to_process"

    pptx_files = glob.glob(f"{directory_to_process}/*.pptx")

    # Process each .pptx file
    for file_path in pptx_files:
        extract_text_to_file(file_path)
